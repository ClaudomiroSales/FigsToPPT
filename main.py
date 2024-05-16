import os
from pptx import Presentation
from pptx.util import Inches

# Define the folder containing the PNG files
folder_path = r'C:\Users\claud\OneDrive\Imagens\Capturas de tela'
# Create a presentation object
presentation = Presentation()

# Iterate over each PNG file in the folder
for file_name in sorted(os.listdir(folder_path)):
    if file_name.lower().endswith('.png'):
        # Create a new slide with a title and content layout
        slide_layout = presentation.slide_layouts[5]  # Using a blank layout
        slide = presentation.slides.add_slide(slide_layout)

        # Define the full path to the image
        image_path = os.path.join(folder_path, file_name)
        
        # Add the image to the slide
        left = Inches(1)
        top = Inches(1)
        height = Inches(5.5)  # Adjust the height as necessary
        slide.shapes.add_picture(image_path, left, top, height=height)

# Save the presentation
presentation.save('output_presentation.pptx')

print("Presentation created successfully!")
