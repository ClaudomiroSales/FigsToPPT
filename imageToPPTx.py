import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# Define the folder containing the PNG files
folder_path = r'C:\Users\claud\OneDrive\Imagens\Capturas de tela'

# Get the dimensions of the first image
first_image_path = os.path.join(folder_path, sorted(os.listdir(folder_path))[0])
with Image.open(first_image_path) as img:
    img_width, img_height = img.size

# Convert pixels to inches (assuming 96 dpi)
dpi = 96
width_in_inches = img_width / dpi
height_in_inches = img_height / dpi

# Create a presentation object and set slide size
presentation = Presentation()
presentation.slide_width = Inches(width_in_inches)
presentation.slide_height = Inches(height_in_inches)

# Iterate over each PNG file in the folder
for file_name in sorted(os.listdir(folder_path)):
    if file_name.lower().endswith('.png'):
        # Create a new slide with a blank layout
        slide_layout = presentation.slide_layouts[6]  # Using a blank layout
        slide = presentation.slides.add_slide(slide_layout)

        # Define the full path to the image
        image_path = os.path.join(folder_path, file_name)
        
        # Add the image to the slide
        left = 0
        top = 0
        slide.shapes.add_picture(image_path, left, top, width=Inches(width_in_inches), height=Inches(height_in_inches))

# Save the presentation
presentation.save('output_presentation.pptx')

print("Presentation created successfully!")
